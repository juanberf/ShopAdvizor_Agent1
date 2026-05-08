import anthropic
import json
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from config.settings import ANTHROPIC_API_KEY, ANTHROPIC_MODEL, OUTPUTS_DIR, LOGO_PATH
from models.session import SessionConfig
from tools.pdf_reader import read_pdf_text


# ── Colores corporativos ──────────────────────────────────────
BLUE          = RGBColor(0x4A, 0x90, 0xD9)
DARK_BLUE     = RGBColor(0x1A, 0x5A, 0xA0)
ORANGE        = RGBColor(0xF5, 0xA6, 0x23)
DARK          = RGBColor(0x2C, 0x3E, 0x50)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
MID_GRAY      = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY    = RGBColor(0xF5, 0xF7, 0xFA)
GREEN         = RGBColor(0x27, 0xAE, 0x60)
RED           = RGBColor(0xE7, 0x4C, 0x3C)
YELLOW        = RGBColor(0xF5, 0xA6, 0x23)
LIGHT_BLUE_BG = RGBColor(0xEB, 0xF3, 0xF9)
GREEN_BG      = RGBColor(0xE8, 0xF5, 0xE9)
RED_BG        = RGBColor(0xFF, 0xF5, 0xF5)
YELLOW_BG     = RGBColor(0xFF, 0xFB, 0xF0)

# ── Layout constants ──────────────────────────────────────────
_SW = 13.33   # slide width  (inches)
_SH = 7.5     # slide height (inches)
_MX = 0.45    # side margin
_CT = 1.22    # content top  (below header + separator)
_CB = 7.08    # content bottom (above footer)


# ── Drawing helpers (also available in Claude's exec namespace) ─

def _r(slide, x, y, w, h, fill=None, lc=None, lw=1.0):
    """Rectangle. All coords in inches. fill/lc are RGBColor."""
    shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if lc is not None:
        shp.line.color.rgb = lc
        shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    return shp


def _t(slide, x, y, w, h, text, pt=12, bold=False, color=None,
       italic=False, align=None, wrap=True):
    """Textbox. All coords in inches."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(pt)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color if color is not None else DARK
    if align:
        p.alignment = align
    return tf


def _logo(slide, x, y, w=1.85):
    """Place logo image. Coords in inches."""
    if LOGO_PATH.exists():
        try:
            slide.shapes.add_picture(str(LOGO_PATH), Inches(x), Inches(y), width=Inches(w))
        except Exception:
            pass


def _header(slide, title, subtitle=None):
    """Standard blue header: bar 0→1.05", orange line 1.05→1.10", logo top-right."""
    _r(slide, 0, 0, _SW, 1.05, fill=BLUE)
    _r(slide, 0, 1.05, _SW, 0.05, fill=ORANGE)
    _t(slide, _MX, 0.14, 10.3, 0.52, title, pt=20, bold=True, color=WHITE)
    if subtitle:
        _t(slide, _MX, 0.66, 10.3, 0.36, subtitle, pt=11,
           color=RGBColor(0xBB, 0xCE, 0xE4))
    # logo pastilla blanca
    _r(slide, 11.18, 0.16, 1.97, 0.73, fill=WHITE)
    _logo(slide, 11.22, 0.20, w=1.85)


def _footer(slide, text="Shopadvizor — Informe Confidencial"):
    _t(slide, _MX, 7.18, _SW - 2 * _MX, 0.28,
       text, pt=9, color=MID_GRAY, align=PP_ALIGN.RIGHT)


CODE_SYSTEM_PROMPT = """Eres un diseñador experto de presentaciones corporativas para Shopadvizor.
Generas código Python con python-pptx que crea presentaciones profesionales, modernas y visualmente impactantes.

══════════════════════════════════════════════════════════════
VARIABLES DISPONIBLES EN EL NAMESPACE (no importes nada)
══════════════════════════════════════════════════════════════
  Presentation, Inches, Pt, RGBColor, PP_ALIGN   ← python-pptx
  os

  Colores:
    BLUE=#4A90D9  DARK_BLUE=#1A5AA0  ORANGE=#F5A623  DARK=#2C3E50
    WHITE  MID_GRAY  LIGHT_GRAY  GREEN=#27AE60  RED=#E74C3C  YELLOW=#F5A623
    LIGHT_BLUE_BG  GREEN_BG  RED_BG  YELLOW_BG

  Constantes de layout (en pulgadas):
    _SW=13.33  _SH=7.5  _MX=0.45  _CT=1.22  _CB=7.08

  Helpers (úsalos siempre — reducen drásticamente el código):
    _r(slide, x, y, w, h, fill=None, lc=None, lw=1.0)
      → rectángulo; coords en pulgadas; fill y lc son RGBColor

    _t(slide, x, y, w, h, text, pt=12, bold=False, color=None,
       italic=False, align=None, wrap=True)
      → textbox; coords en pulgadas; devuelve text_frame

    _logo(slide, x, y, w=1.85)
      → añade el logo PNG; coords en pulgadas

    _header(slide, title, subtitle=None)
      → cabecera estándar: barra BLUE, línea ORANGE, logo arriba-derecha

    _footer(slide, text="...")
      → texto gris abajo-derecha

  Rutas:
    OUTPUT_PATH  → str — ruta donde guardar el .pptx
    LOGO_PATH    → Path — ruta al logo (usa _logo() para añadirlo)

══════════════════════════════════════════════════════════════
DIMENSIONES — CRÍTICO: nunca superar
══════════════════════════════════════════════════════════════
  Slide: 13.33" × 7.5"
  Zona de contenido bajo _header: y ∈ [1.22, 7.08], x ∈ [0.45, 12.88]
  Regla: x+w ≤ 13.33 y y+h ≤ 7.5 en TODOS los elementos

══════════════════════════════════════════════════════════════
API python-pptx — REFERENCIA MÍNIMA NECESARIA
══════════════════════════════════════════════════════════════
  prs = Presentation()
  prs.slide_width  = Inches(13.33)
  prs.slide_height = Inches(7.5)

  slide = prs.slides.add_slide(prs.slide_layouts[6])   ← SIEMPRE layout[6]

  # Fondo sólido:
  slide.background.fill.solid()
  slide.background.fill.fore_color.rgb = DARK_BLUE

  # Añadir párrafo a un tf existente:
  p2 = tf.add_paragraph()
  p2.text = "Línea 2"
  p2.font.size = Pt(11)
  p2.font.color.rgb = MID_GRAY
  p2.space_before = Pt(5)

  # Tabla:
  tbl = slide.shapes.add_table(n_rows, n_cols,
          Inches(x), Inches(y), Inches(w), Inches(h)).table
  cell = tbl.cell(row, col)
  cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
  cp = cell.text_frame.paragraphs[0]
  cp.text = "Val"; cp.font.size = Pt(10); cp.font.color.rgb = WHITE

  prs.save(OUTPUT_PATH)   ← OBLIGATORIO al final

══════════════════════════════════════════════════════════════
DISEÑO — LIBERTAD CREATIVA CON IDENTIDAD CORPORATIVA
══════════════════════════════════════════════════════════════
  • Usa _r y _t para todo — son 1 línea en lugar de 5-8
  • Usa _header(slide, "Título") en todas las slides de contenido
  • Usa _footer(slide) en todas las slides de contenido
  • Portada: fondo DARK_BLUE, franja ORANGE izquierda (x=0, w=0.35, h=7.5),
    título 36pt+ blanco, panel KPI lateral derecho (x≈9.8), logo pastilla blanca abajo
  • Slides de contenido: _header → zona y∈[_CT, _CB] → _footer
  • Closing: fondo DARK_BLUE, franja ORANGE, puntos grandes blancos, logo abajo-derecha
  • Valores numéricos clave: ≥22pt, color según estado (GREEN positivo, RED crítico, YELLOW atención)
  • Usa ORANGE para separadores y franjas de acento
  • Prioriza datos y números sobre texto largo — una idea por slide

══════════════════════════════════════════════════════════════
REGLAS DE CONTENIDO — OBLIGATORIAS
══════════════════════════════════════════════════════════════
  1. Slide 1: portada con título, subtítulo, meta y KPIs flash
  2. Slide 2: resumen ejecutivo con párrafo de síntesis y puntos clave
  3. Penúltima: recomendaciones accionables (≥4, con cifras que las justifiquen)
  4. Última: conclusiones estratégicas
  5. NUNCA inventes datos — usa solo lo del JSON
  6. NUNCA posiciones exactas en rankings (#1, posición 3...)
     → "parte alta de la categoría", "por encima de la media", "en el tramo superior"
  7. Sin nombres de competidores
  8. Mínimo 5 slides · Máximo 7
  9. El código completo NO debe superar 280 líneas — usa bucles y comprensiones de lista

══════════════════════════════════════════════════════════════
FORMATO DE RESPUESTA — CRÍTICO
══════════════════════════════════════════════════════════════
  Solo código Python ejecutable. Sin markdown. Sin comentarios. Sin texto extra.
  Primera línea: prs = Presentation()
  Última línea:  prs.save(OUTPUT_PATH)"""


def run(
    session: SessionConfig,
    validated_data: dict,
    raw_data: dict,
    analysis_data: dict,
    alerts_data: dict,
    insights_text: str,
    user_request: str = "",
    pdf_path: Path = None,
) -> Path:
    """
    Genera un PowerPoint: Claude diseña y escribe el código python-pptx,
    que se ejecuta directamente para producir el archivo final.
    """
    print(f"[SA6] Iniciando generación de PowerPoint para sesión {session.session_id}")

    client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)

    pdf_text = None
    if pdf_path and Path(pdf_path).exists():
        try:
            pdf_text = read_pdf_text(pdf_path)
            print(f"[SA6] PDF leído: {len(pdf_text)} chars")
        except Exception as e:
            print(f"[SA6] No se pudo leer el PDF: {e}")

    campaign_data = _prepare_pptx_data(
        validated_data, analysis_data, alerts_data, insights_text, pdf_text
    )

    session_outputs = OUTPUTS_DIR / session.session_id
    session_outputs.mkdir(parents=True, exist_ok=True)
    pptx_path = session_outputs / "presentacion.pptx"

    lang = "ESPAÑOL" if session.language == "es" else "FRANÇAIS"

    print("[SA6] Generando código de presentación con Claude...")
    code = _generate_pptx_code(client, campaign_data, lang, session, user_request, pptx_path)

    print("[SA6] Ejecutando código generado...")
    _execute_pptx_code(code, pptx_path)

    print(f"[SA6] PowerPoint generado en {pptx_path}")
    return pptx_path


def _generate_pptx_code(
    client: anthropic.Anthropic,
    campaign_data: dict,
    lang: str,
    session: SessionConfig,
    user_request: str,
    output_path: Path,
) -> str:
    """Llama a Claude para que genere el código python-pptx completo."""
    message = client.messages.create(
        model=ANTHROPIC_MODEL,
        max_tokens=8192,
        system=CODE_SYSTEM_PROMPT,
        messages=[{
            "role": "user",
            "content": f"""Idioma: {lang}
Tono: {session.tone}
Segmentos de foco: {', '.join(session.focus_segments) if session.focus_segments else 'Todos'}
Petición: {user_request or 'Presentación ejecutiva completa de la campaña'}

Valores disponibles en el namespace:
  OUTPUT_PATH = "{output_path}"
  LOGO_PATH   = Path("{LOGO_PATH}")

DATOS DE LA CAMPAÑA:
{json.dumps(campaign_data, ensure_ascii=False, indent=2)}

Genera el código ahora. Usa _r, _t, _header, _footer para mantenerlo conciso."""
        }],
    )

    code = message.content[0].text.strip()
    if "```python" in code:
        code = code.split("```python")[1].split("```")[0].strip()
    elif "```" in code:
        code = code.split("```")[1].split("```")[0].strip()

    return code


def _execute_pptx_code(code: str, output_path: Path) -> None:
    """Ejecuta el código python-pptx generado por Claude en un namespace controlado."""
    namespace = {
        "__builtins__": __builtins__,
        # python-pptx
        "Presentation":   Presentation,
        "Inches":         Inches,
        "Pt":             Pt,
        "Emu":            Emu,
        "RGBColor":       RGBColor,
        "PP_ALIGN":       PP_ALIGN,
        # stdlib
        "os":             os,
        # layout constants
        "_SW": _SW, "_SH": _SH, "_MX": _MX, "_CT": _CT, "_CB": _CB,
        # drawing helpers
        "_r":      _r,
        "_t":      _t,
        "_logo":   _logo,
        "_header": _header,
        "_footer": _footer,
        # colores corporativos
        "BLUE":           BLUE,
        "DARK_BLUE":      DARK_BLUE,
        "ORANGE":         ORANGE,
        "DARK":           DARK,
        "WHITE":          WHITE,
        "MID_GRAY":       MID_GRAY,
        "LIGHT_GRAY":     LIGHT_GRAY,
        "GREEN":          GREEN,
        "RED":            RED,
        "YELLOW":         YELLOW,
        "LIGHT_BLUE_BG":  LIGHT_BLUE_BG,
        "GREEN_BG":       GREEN_BG,
        "RED_BG":         RED_BG,
        "YELLOW_BG":      YELLOW_BG,
        # rutas
        "LOGO_PATH":      LOGO_PATH,
        "OUTPUT_PATH":    str(output_path),
    }
    exec(code, namespace)  # noqa: S102
    if not output_path.exists():
        raise RuntimeError("El código generado no produjo el archivo PowerPoint esperado")


def _prepare_pptx_data(
    validated_data: dict,
    analysis_data: dict,
    alerts_data: dict,
    insights_text: str,
    pdf_text: str = None,
) -> dict:
    """
    Prepara datos resumidos para el prompt de Claude.
    No incluye preguntas pre/post raw — SA4 ya las ha analizado en analysis_data.
    """
    metadata    = validated_data.get("campaign_metadata", {})
    kpis        = validated_data.get("kpis", {})
    attrs       = validated_data.get("product_attributes", {})
    qualitative = validated_data.get("qualitative", {})
    competitive = validated_data.get("competitive", {})
    pre_eval    = validated_data.get("pre_evaluation", {})
    post_eval   = validated_data.get("post_evaluation", {})

    alert_levels = {}
    for alert in alerts_data.get("alerts", []):
        kpi_name = alert.get("kpi", "")
        if not kpi_name.startswith("ranking_") and kpi_name != "pre_post_purchase_intent":
            alert_levels[kpi_name] = alert.get("level")

    kpis_enriched = {}
    for name, data in kpis.items():
        if isinstance(data, dict):
            kpis_enriched[name] = {**data, "alert_level": alert_levels.get(name, "ok")}

    # Resumen cualitativo compacto
    qualitative_summary = {
        "sentiment":    qualitative.get("sentiment", {}),
        "strengths":    qualitative.get("strengths", [])[:4],
        "improvements": qualitative.get("improvements", [])[:3],
    }

    # Top atributos de producto (los 4 mejores y los 2 peores)
    all_attrs = {}
    for group in attrs.values() if isinstance(attrs, dict) else []:
        if isinstance(group, dict):
            all_attrs.update(group)
    sorted_attrs = sorted(all_attrs.items(),
                          key=lambda x: x[1].get("mean", 0) if isinstance(x[1], dict) else 0,
                          reverse=True)
    attrs_summary = {k: v for k, v in sorted_attrs[:4] + sorted_attrs[-2:]}

    result = {
        "campaign_metadata": metadata,
        "kpis": kpis_enriched,
        "segments": {
            "pre": pre_eval.get("segments", []),
            "post": post_eval.get("segments", []),
        },
        "qualitative": qualitative_summary,
        "product_attributes_highlights": attrs_summary,
        "competitive": competitive,
        "alerts_summary": alerts_data.get("summary", {}),
        "alerts": alerts_data.get("alerts", [])[:10],
        "analysis": analysis_data,
        "insights": insights_text[:1500] if insights_text else "",
    }

    if pdf_text:
        result["pdf_context"] = pdf_text[:2000]

    return result
